import pptx


class Reader():
    def __init__(self, charLimit=500):
        self.charLimit = charLimit

    def parse(self, filePath):
        try:
            fileContent = ""
            pres = pptx.Presentation(filePath)
            for slide in pres.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        fileContent += shape.text
            if len(fileContent) > self.charLimit:
                fileContent = fileContent[:self.charLimit]
            return fileContent
        except Exception:
            return None
